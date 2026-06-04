from __future__ import annotations

import os
import platform
import subprocess
import re
import uuid
from pathlib import Path
from typing import BinaryIO

from db import DATA_DIR, execute_many, write_transaction
from services.auth_service import require_login
from services.pdf_extraction_service import extract_pdf_pages as extract_pdf_pages_from_pdf

UPLOAD_DIR = DATA_DIR / "uploads"
PAGE_IMAGE_DIR = DATA_DIR / "page_images"
PAGE_SOURCE_DIR = DATA_DIR / "page_sources"


def import_deck(uploaded_file: BinaryIO, *, subject: str, title: str) -> int:
    user = require_login()
    data = bytes(uploaded_file.getbuffer() if hasattr(uploaded_file, "getbuffer") else uploaded_file.read())
    upload_size = len(data)
    if not _reserve_upload_capacity(user.id, upload_size):
        raise RuntimeError("上传失败：已超过当前账户的上传容量配额。请联系管理员扩容或先删除旧资料。")
    original_name = getattr(uploaded_file, "name", "")
    suffix = Path(original_name).suffix.lower()
    if suffix == ".pptx":
        return import_pptx(uploaded_file, subject=subject, title=title, file_bytes=data)
    if suffix == ".pdf":
        return import_pdf(uploaded_file, subject=subject, title=title, file_bytes=data)
    raise RuntimeError("仅支持 PPTX 或 PDF 文件。")


def import_pptx(uploaded_file: BinaryIO, *, subject: str, title: str, file_bytes: bytes | None = None) -> int:
    saved_path = save_uploaded_deck(uploaded_file, file_bytes=file_bytes)
    slides = extract_pptx_slides(saved_path)
    image_paths = render_deck_page_images(saved_path)
    return _save_deck_records(saved_path, slides, image_paths, subject=subject, title=title)


def import_pdf(uploaded_file: BinaryIO, *, subject: str, title: str, file_bytes: bytes | None = None) -> int:
    saved_path = save_uploaded_deck(uploaded_file, file_bytes=file_bytes)
    slides = extract_pdf_pages(saved_path)
    image_paths = render_deck_page_images(saved_path)
    return _save_deck_records(saved_path, slides, image_paths, subject=subject, title=title)


def _save_deck_records(
    saved_path: Path,
    slides: list[dict[str, str | int]],
    image_paths: dict[int, Path],
    *,
    subject: str,
    title: str,
) -> int:
    user = require_login()
    deck_title = title.strip() or saved_path.stem
    with write_transaction() as conn:
        if not _has_upload_capacity(conn, user.id, saved_path.stat().st_size):
            try:
                saved_path.unlink()
            except OSError:
                pass
            for image_path in image_paths.values():
                try:
                    image_path.unlink()
                except OSError:
                    pass
            raise RuntimeError("上传失败：已超过当前账户的上传容量配额。请联系管理员扩容或先删除旧资料。")
        cursor = conn.execute(
            """
            INSERT INTO ppt_decks (user_id, filename, title, subject, file_path, slide_count)
            VALUES (?, ?, ?, ?, ?, ?)
            """,
            (user.id, saved_path.name, deck_title, subject.strip(), str(saved_path), len(slides)),
        )
        deck_id = int(cursor.lastrowid)
        conn.executemany(
            """
            INSERT INTO ppt_slides (user_id, deck_id, slide_number, title, slide_text, notes, image_path)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (
                (
                    user.id,
                    deck_id,
                    slide["slide_number"],
                    slide["title"],
                    slide["slide_text"],
                    slide["notes"],
                    str(image_paths.get(int(slide["slide_number"]), "")),
                )
                for slide in slides
            ),
        )
    return deck_id


def save_uploaded_deck(uploaded_file: BinaryIO, *, file_bytes: bytes | None = None) -> Path:
    user = require_login()
    user_upload_dir = UPLOAD_DIR / f"user_{user.id}"
    user_upload_dir.mkdir(parents=True, exist_ok=True)
    original_name = getattr(uploaded_file, "name", "uploaded")
    suffix = Path(original_name).suffix.lower() or ".pptx"
    safe_stem = _safe_filename(Path(original_name).stem) or "deck"
    target = user_upload_dir / f"{uuid.uuid4().hex}_{safe_stem}{suffix}"
    temp_target = target.with_name(f".{target.name}.tmp")

    data = file_bytes if file_bytes is not None else bytes(uploaded_file.getbuffer() if hasattr(uploaded_file, "getbuffer") else uploaded_file.read())
    temp_target.write_bytes(data)
    temp_target.replace(target)
    return target


def save_page_source_file(uploaded_file: BinaryIO, *, file_bytes: bytes | None = None) -> Path:
    user = require_login()
    original_name = getattr(uploaded_file, "name", "uploaded")
    suffix = Path(original_name).suffix.lower()
    if suffix not in {".pptx", ".pdf"}:
        raise RuntimeError("仅支持 PPTX 或 PDF 文件。")

    user_source_dir = PAGE_SOURCE_DIR / f"user_{user.id}"
    user_source_dir.mkdir(parents=True, exist_ok=True)
    safe_stem = _safe_filename(Path(original_name).stem) or "source"
    target = user_source_dir / f"{uuid.uuid4().hex}_{safe_stem}{suffix}"
    temp_target = target.with_name(f".{target.name}.tmp")

    data = file_bytes if file_bytes is not None else bytes(uploaded_file.getbuffer() if hasattr(uploaded_file, "getbuffer") else uploaded_file.read())
    temp_target.write_bytes(data)
    temp_target.replace(target)
    return target


def extract_pptx_slides(path: Path) -> list[dict[str, str | int]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("未安装 python-pptx，请先运行 pip install -r requirements.txt。") from exc

    presentation = Presentation(path)
    slides: list[dict[str, str | int]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = _extract_title(slide)
        text_lines = _extract_text_lines(slide)
        slides.append(
            {
                "slide_number": index,
                "title": title,
                "slide_text": "\n".join(text_lines),
                "notes": "",
            }
        )
    return slides


def extract_pdf_pages(path: Path, *, method: str = "local") -> list[dict[str, str | int]]:
    return extract_pdf_pages_from_pdf(path, method=method)


def extract_source_pages(path: Path, *, method: str = "local") -> list[dict[str, str | int]]:
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        pages = extract_pdf_pages(path, method=method)
    elif suffix == ".pptx":
        pages = extract_pptx_slides(path)
    else:
        raise RuntimeError("仅支持解析 PPTX 或 PDF 文件。")

    image_paths = render_deck_page_images(path)
    result: list[dict[str, str | int]] = []
    for page in sorted(pages, key=lambda item: int(item["slide_number"])):
        slide_number = int(page["slide_number"])
        result.append(
            {
                "slide_number": slide_number,
                "title": str(page.get("title") or f"第 {slide_number} 页"),
                "slide_text": str(page.get("slide_text") or ""),
                "notes": str(page.get("notes") or ""),
                "image_path": str(image_paths.get(slide_number, "")),
            }
        )
    return result


def render_deck_page_images(path: Path) -> dict[int, Path]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return render_pdf_page_images(path)
    if suffix == ".pptx":
        return render_pptx_page_images(path)
    raise RuntimeError("仅支持渲染 PDF 或 PPTX。")


def render_pdf_page_images(path: Path) -> dict[int, Path]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("未安装 PyMuPDF，请先运行 pip install -r requirements.txt。") from exc

    target_dir = _page_image_target_dir(path)
    target_dir.mkdir(parents=True, exist_ok=True)
    document = fitz.open(path)
    result: dict[int, Path] = {}
    matrix = fitz.Matrix(2, 2)
    for index, page in enumerate(document, start=1):
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        target = target_dir / f"page_{index:03d}.png"
        pixmap.save(target)
        result[index] = target
    document.close()
    return result


def render_pptx_page_images(path: Path) -> dict[int, Path]:
    if platform.system() == "Windows":
        return _render_pptx_page_images_windows(path)
    else:
        return _render_pptx_page_images_linux(path)


def _render_pptx_page_images_windows(path: Path) -> dict[int, Path]:
    target_dir = _page_image_target_dir(path)
    target_dir.mkdir(parents=True, exist_ok=True)

    script = f"""
$ErrorActionPreference = 'Stop'
$pptPath = {str(path.resolve())!r}
$outDir = {str(target_dir.resolve())!r}
$app = New-Object -ComObject PowerPoint.Application
try {{
    $presentation = $app.Presentations.Open($pptPath, $false, $false, $false)
    try {{
        $presentation.Export($outDir, 'PNG', 1440, 810)
    }} finally {{
        $presentation.Close()
    }}
}} finally {{
    $app.Quit()
    [System.Runtime.InteropServices.Marshal]::ReleaseComObject($app) | Out-Null
}}
"""
    completed = subprocess.run(
        ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", script],
        capture_output=True,
        text=True,
        timeout=180,
    )
    if completed.returncode != 0:
        raise RuntimeError(f"PowerPoint 导出页面图片失败：{completed.stderr or completed.stdout}")

    result: dict[int, Path] = {}
    for image in target_dir.glob("*.PNG"):
        index = _slide_image_index(image.name)
        if index:
            normalized = target_dir / f"page_{index:03d}.png"
            if image != normalized:
                image.replace(normalized)
            result[index] = normalized
    for image in target_dir.glob("*.png"):
        index = _slide_image_index(image.name)
        if index:
            result[index] = image
    if not result:
        raise RuntimeError("PowerPoint 导出完成，但未找到页面图片。")
    return result


def _render_pptx_page_images_linux(path: Path) -> dict[int, Path]:
    result: dict[int, Path] = {}
    # Try LibreOffice first
    lo_result = _render_pptx_with_libreoffice(path)
    if lo_result:
        return lo_result
    # Fallback: try python-pptx + PIL to render each slide
    return _render_pptx_with_pptx2img(path)


def _render_pptx_with_libreoffice(path: Path) -> dict[int, Path] | None:
    target_dir = _page_image_target_dir(path)
    target_dir.mkdir(parents=True, exist_ok=True)
    try:
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "png",
                "--outdir",
                str(target_dir),
                str(path),
            ],
            capture_output=True,
            text=True,
            timeout=180,
        )
    except FileNotFoundError:
        return None  # LibreOffice not installed
    if result.returncode != 0:
        return None  # LibreOffice failed

    result_map: dict[int, Path] = {}
    for image in sorted(target_dir.glob("*.png")):
        index = _slide_image_index(image.name)
        if index:
            normalized = target_dir / f"page_{index:03d}.png"
            if image != normalized:
                image.replace(normalized)
            result_map[index] = normalized
    return result_map if result_map else None


def _render_pptx_with_pptx2img(path: Path) -> dict[int, Path]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from PIL import Image
        import io
    except ImportError as exc:
        raise RuntimeError(
            "在 Linux 上渲染 PPTX 页面图片需要安装依赖：pip install python-pptx Pillow\n"
            "或安装 LibreOffice 后将其添加到 PATH。"
        ) from exc

    target_dir = _page_image_target_dir(path)
    target_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(path)

    result: dict[int, Path] = {}
    for index, slide in enumerate(presentation.slides, start=1):
        # Get slide dimensions
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        # Create a blank image
        width_px = int(slide_width / 914400) * 2  # Convert EMUs to pixels at 2x scale
        height_px = int(slide_height / 914400) * 2
        img = Image.new("RGB", (width_px, height_px), (255, 255, 255))

        # Try to render using pptx's shape positioning
        # Note: python-pptx cannot directly render to image,
        # so we create a visual representation from shapes
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for shape in slide.shapes:
            if shape.has_text_frame:
                _draw_text_shape(img, shape, slide_width, slide_height)
            elif hasattr(shape, "image"):
                _draw_picture_shape(img, shape, slide_width, slide_height)

        target = target_dir / f"page_{index:03d}.png"
        img.save(target)
        result[index] = target

    if not result:
        raise RuntimeError("PPTX 页面渲染失败，未生成任何图片。")
    return result


def _draw_text_shape(img: Image.Image, shape, slide_width: int, slide_height: int) -> None:
    try:
        from PIL import ImageDraw, ImageFont
        left = int(shape.left / 914400) * 2
        top = int(shape.top / 914400) * 2
        width = int(shape.width / 914400) * 2
        height = int(shape.height / 914400) * 2
        draw = ImageDraw.Draw(img)
        text = shape.text_frame.text
        if text.strip():
            # Simple text rendering - just draw text
            for line in text.split("\n")[:20]:  # Limit lines
                if line.strip():
                    draw.text((left + 5, top + 5), line[:100], fill=(0, 0, 0))
                    top += 20
    except Exception:
        pass


def _draw_picture_shape(img: Image.Image, shape, slide_width: int, slide_height: int) -> None:
    try:
        left = int(shape.left / 914400) * 2
        top = int(shape.top / 914400) * 2
        width = int(shape.width / 914400) * 2
        height = int(shape.height / 914400) * 2
        img.paste((200, 200, 200), (left, top, left + width, top + height))
    except Exception:
        pass


def render_missing_page_images(deck: dict, slides: list[dict]) -> dict[int, Path]:
    path = Path(deck["file_path"])
    image_paths = render_deck_page_images(path)
    execute_many(
        "UPDATE ppt_slides SET image_path = ? WHERE id = ? AND user_id = ?",
        (
            (str(image_path), slide["id"], slide["user_id"])
            for slide in slides
            if (image_path := image_paths.get(int(slide["slide_number"])))
        ),
    )
    return image_paths


def apply_source_page_to_deck(
    deck: dict,
    source_page: dict,
    *,
    target_slide_number: int,
    mode: str,
) -> dict[str, int]:
    user = require_login()
    deck_id = int(deck["id"])
    target = int(target_slide_number)
    normalized_mode = mode.strip().lower()
    if normalized_mode not in {"insert", "replace"}:
        raise RuntimeError("页面操作仅支持 insert 或 replace。")

    page = _normalize_source_page(source_page)
    with write_transaction() as conn:
        max_slide_number = _max_slide_number(conn, user.id, deck_id)
        if normalized_mode == "insert":
            if target < 1 or target > max_slide_number + 1:
                raise RuntimeError(f"插入位置必须在 1 到 {max_slide_number + 1} 之间。")
            _shift_slide_numbers_up(conn, user.id, deck_id, target)
            cursor = conn.execute(
                """
                INSERT INTO ppt_slides (user_id, deck_id, slide_number, title, slide_text, notes, image_path)
                VALUES (?, ?, ?, ?, ?, ?, ?)
                """,
                (user.id, deck_id, target, page["title"], page["slide_text"], page["notes"], page["image_path"]),
            )
            conn.execute(
                "UPDATE ppt_decks SET slide_count = ? WHERE id = ? AND user_id = ?",
                (max_slide_number + 1, deck_id, user.id),
            )
            return {"inserted": 1, "replaced": 0, "slide_id": int(cursor.lastrowid), "target_slide_number": target}

        if target < 1 or target > max_slide_number:
            raise RuntimeError(f"替换位置必须在 1 到 {max_slide_number} 之间。")
        cursor = conn.execute(
            """
            UPDATE ppt_slides
            SET title = ?, slide_text = ?, notes = ?, image_path = ?
            WHERE user_id = ? AND deck_id = ? AND slide_number = ?
            """,
            (page["title"], page["slide_text"], page["notes"], page["image_path"], user.id, deck_id, target),
        )
        if cursor.rowcount == 0:
            raise RuntimeError(f"未找到第 {target} 页，无法替换。")
        return {"inserted": 0, "replaced": 1, "target_slide_number": target}


def delete_deck_page(deck: dict, *, target_slide_number: int) -> dict[str, int]:
    user = require_login()
    deck_id = int(deck["id"])
    target = int(target_slide_number)
    with write_transaction() as conn:
        max_slide_number = _max_slide_number(conn, user.id, deck_id)
        if target < 1 or target > max_slide_number:
            raise RuntimeError(f"删除位置必须在 1 到 {max_slide_number} 之间。")
        row = conn.execute(
            """
            SELECT id
            FROM ppt_slides
            WHERE user_id = ? AND deck_id = ? AND slide_number = ?
            """,
            (user.id, deck_id, target),
        ).fetchone()
        if row is None:
            raise RuntimeError(f"未找到第 {target} 页，无法删除。")
        slide_id = int(row[0])
        conn.execute("DELETE FROM slide_questions WHERE user_id = ? AND slide_id = ?", (user.id, slide_id))
        conn.execute("DELETE FROM slide_explanations WHERE user_id = ? AND slide_id = ?", (user.id, slide_id))
        conn.execute("DELETE FROM ppt_slides WHERE user_id = ? AND id = ?", (user.id, slide_id))
        conn.execute(
            "DELETE FROM ppt_study_asset_pages WHERE user_id = ? AND deck_id = ? AND slide_number = ?",
            (user.id, deck_id, target),
        )
        _shift_slide_numbers_down(conn, user.id, deck_id, target)
        conn.execute(
            """
            UPDATE ppt_sections
            SET
                start_slide = CASE WHEN start_slide > ? THEN start_slide - 1 ELSE start_slide END,
                end_slide = CASE WHEN end_slide >= ? THEN end_slide - 1 ELSE end_slide END,
                updated_at = datetime('now', 'localtime')
            WHERE user_id = ? AND deck_id = ?
            """,
            (target, target, user.id, deck_id),
        )
        conn.execute(
            "DELETE FROM ppt_sections WHERE user_id = ? AND deck_id = ? AND end_slide < start_slide",
            (user.id, deck_id),
        )
        conn.execute(
            "UPDATE ppt_decks SET slide_count = ? WHERE id = ? AND user_id = ?",
            (max(0, max_slide_number - 1), deck_id, user.id),
        )
    return {"deleted": 1, "target_slide_number": target}


def _normalize_source_page(source_page: dict) -> dict[str, str]:
    slide_number = int(source_page.get("slide_number") or 0)
    return {
        "title": str(source_page.get("title") or f"第 {slide_number} 页"),
        "slide_text": str(source_page.get("slide_text") or ""),
        "notes": str(source_page.get("notes") or ""),
        "image_path": str(source_page.get("image_path") or ""),
    }


def _max_slide_number(conn, user_id: int, deck_id: int) -> int:
    row = conn.execute(
        """
        SELECT COALESCE(MAX(slide_number), 0)
        FROM ppt_slides
        WHERE user_id = ? AND deck_id = ?
        """,
        (user_id, deck_id),
    ).fetchone()
    return int(row[0] or 0) if row else 0


def _shift_slide_numbers_up(conn, user_id: int, deck_id: int, start_slide_number: int) -> None:
    conn.execute(
        """
        UPDATE ppt_slides
        SET slide_number = -slide_number
        WHERE user_id = ? AND deck_id = ? AND slide_number >= ?
        """,
        (user_id, deck_id, start_slide_number),
    )
    conn.execute(
        """
        UPDATE ppt_slides
        SET slide_number = -slide_number + 1
        WHERE user_id = ? AND deck_id = ? AND slide_number < 0
        """,
        (user_id, deck_id),
    )
    conn.execute(
        """
        UPDATE ppt_study_asset_pages
        SET slide_number = slide_number + 1
        WHERE user_id = ? AND deck_id = ? AND slide_number >= ?
        """,
        (user_id, deck_id, start_slide_number),
    )
    conn.execute(
        """
        UPDATE ppt_sections
        SET
            start_slide = CASE WHEN start_slide >= ? THEN start_slide + 1 ELSE start_slide END,
            end_slide = CASE WHEN end_slide >= ? THEN end_slide + 1 ELSE end_slide END,
            updated_at = datetime('now', 'localtime')
        WHERE user_id = ? AND deck_id = ?
        """,
        (start_slide_number, start_slide_number, user_id, deck_id),
    )


def _shift_slide_numbers_down(conn, user_id: int, deck_id: int, deleted_slide_number: int) -> None:
    conn.execute(
        """
        UPDATE ppt_slides
        SET slide_number = -slide_number
        WHERE user_id = ? AND deck_id = ? AND slide_number > ?
        """,
        (user_id, deck_id, deleted_slide_number),
    )
    conn.execute(
        """
        UPDATE ppt_slides
        SET slide_number = -slide_number - 1
        WHERE user_id = ? AND deck_id = ? AND slide_number < 0
        """,
        (user_id, deck_id),
    )
    conn.execute(
        """
        UPDATE ppt_study_asset_pages
        SET slide_number = slide_number - 1
        WHERE user_id = ? AND deck_id = ? AND slide_number > ?
        """,
        (user_id, deck_id, deleted_slide_number),
    )


def _page_image_target_dir(path: Path) -> Path:
    user_id = _user_id_from_upload_path(path) or require_login().id
    return PAGE_IMAGE_DIR / f"user_{user_id}" / _safe_filename(path.stem)


def _slide_image_index(filename: str) -> int | None:
    match = re.search(r"(\d+)", filename)
    return int(match.group(1)) if match else None


def _extract_title(slide) -> str:
    if slide.shapes.title and getattr(slide.shapes.title, "text", "").strip():
        return slide.shapes.title.text.strip()
    for line in _extract_text_lines(slide):
        if line.strip():
            return line.strip()[:80]
    return "未命名页面"


def _extract_text_lines(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        if not text:
            continue
        for line in text.splitlines():
            clean = line.strip()
            if clean:
                lines.append(clean)
    return lines


def _first_text_line(text: str) -> str:
    for line in text.splitlines():
        clean = line.strip()
        if clean:
            return clean
    return ""


def _safe_filename(value: str) -> str:
    return re.sub(r"[^0-9A-Za-z\u4e00-\u9fff._-]+", "_", value).strip("._-")


def _reserve_upload_capacity(user_id: int, upload_size: int) -> bool:
    return True


def _has_upload_capacity(conn, user_id: int, upload_size: int) -> bool:
    return True


def _user_id_from_upload_path(path: Path) -> int | None:
    for parent in path.parents:
        match = re.fullmatch(r"user_(\d+)", parent.name)
        if match:
            return int(match.group(1))
    return None
